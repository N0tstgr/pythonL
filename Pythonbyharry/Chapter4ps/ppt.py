from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a new PowerPoint presentation
prs = Presentation()

# Define slide titles and content from the handwritten notes (summarized into 15 slides)
slides_content = [
    ("Marketing Management", "Assignment on Cricket Bat"),
    ("Why I Chose This Product", 
     "I chose the cricket bat because cricket is one of the most loved sports in our country. It connects people of all ages and backgrounds. The bat is the most important equipment, and I find its design and making process very interesting."),
    ("Why is it Important?", 
     "A cricket bat is essential for playing cricket. It helps the player/batsman hit the ball, score runs, and win matches. A good quality bat provides comfort, balance, and performance."),
    ("Materials Used in Making Bat", 
     "- Willow wood (Kashmir or England)\n- Good cane and rubber for handle\n- Thread and glue for binding\n- Rubber grip for hold\n- Polish and stickers for finishing"),
    ("Economic Impact", 
     "Cricket bat production helps the economy:\n- Creates employment in woodwork and sports goods industry\n- Encourages local craftsmanship\n- Boosts exports when bats are sold internationally"),
    ("Selling Methods", 
     "Cricket bats can be sold both online and offline:\nOffline - sports shops, wholesalers, schools\nOnline - websites, Instagram, Amazon, Flipkart, Meesho"),
    ("Target Countries for Export", 
     "Target countries where cricket is growing:\nUSA, Canada, UAE, Nepal, Germany"),
    ("Marketing Strategy", 
     "Create brand website or Instagram page.\nCollaborate with online platforms for food market.\nSend samples to international retailers."),
    ("Why Selling Bats is a Good Business", 
     "High demand in cricket-loving nations.\nSchools and colleges organize matches.\nLow investment, high returns.\nEveryone loves cricket."),
    ("Cricket Bat and Youth", 
     "Young people love cricket.\nHigh potential market among youth and schools."),
    ("Business Opportunities", 
     "Start small with local suppliers.\nScale up to international market."),
    ("Export Potential", 
     "International demand is rising.\nMore countries showing interest in cricket."),
    ("Encouraging Local Talent", 
     "Helps local craftsmen and promotes indigenous skills."),
    ("Conclusion", 
     "Cricket bats are not just sports equipment but a strong business opportunity and cultural product."),
    ("Thank You", "Presentation by [Your Name]")
]

# Function to add a slide with title and content
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = content
    # Style
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)
                run.font.name = 'Calibri'

# Create slides
for slide_title, slide_content in slides_content:
    add_slide(slide_title, slide_content)

# Save the presentation
pptx_path = "/mnt/data/Cricket_Bat_Assignment_Presentation.pptx"
prs.save(pptx_path)

pptx_path
